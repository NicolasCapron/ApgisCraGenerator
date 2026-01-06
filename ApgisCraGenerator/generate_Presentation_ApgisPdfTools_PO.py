#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# Génère Presentation_ApgisPdfTools_PO.pptx (français) pour Product Owners
# Dépendance : python-pptx
# Installer : pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt

def add_title_slide(prs, title, subtitle=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_tf = slide.shapes.title
    title_tf.text = title
    if subtitle:
        try:
            subtitle_tf = slide.placeholders[1]
            subtitle_tf.text = subtitle
        except:
            pass
    return slide

def add_bullets_slide(prs, title, bullets, notes=None):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    return slide

def add_two_column_slide(prs, title, left_bullets, right_bullets, notes=None):
    slide_layout = prs.slide_layouts[3]  # Two Content (may vary by template)
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    # left
    left_tf = slide.shapes.placeholders[1].text_frame
    left_tf.clear()
    for i,b in enumerate(left_bullets):
        if i==0:
            p = left_tf.paragraphs[0]
            p.text = b
        else:
            p = left_tf.add_paragraph()
            p.text = b
        p.font.size = Pt(16)
    # right
    right_tf = slide.shapes.placeholders[2].text_frame
    right_tf.clear()
    for i,b in enumerate(right_bullets):
        if i==0:
            p = right_tf.paragraphs[0]
            p.text = b
        else:
            p = right_tf.add_paragraph()
            p.text = b
        p.font.size = Pt(16)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def main():
    prs = Presentation()
    # Slide 1 - Titre
    add_title_slide(prs,
                    "Migration de la solution PDF pour Apgis",
                    "Présentation de la classe ApgisPdfTools (PdfTools SDK) — Public : Product Owners")
    # Slide 2 - Agenda
    add_bullets_slide(prs, "Agenda", [
        "Contexte & historique rapide",
        "Pourquoi on change (bénéfices métier)",
        "Vue fonctionnelle de la solution",
        "Impacts à prévoir (technique, process, licence)",
        "Risques et mitigations",
        "Plan de migration & critères d'acceptation",
        "Décisions à prendre / prochaines étapes"
    ], notes="Présenter l'ordre et le timing : ~30-45 min (présentation + Q&A).")
    # Slide 3 - Historique
    add_bullets_slide(prs, "Historique rapide (contexte)", [
        "Jusqu'ici : iTextSharp pour manipulations PDF (fusion, extraction, génération).",
        "Limites : conformité PDF/A, maintenance, licence/compatibilité.",
        "Décision : adopter PdfTools SDK via la classe ApgisPdfTools.",
        "But métier : garantir pérennité des documents archivés et réduire incidents."
    ], notes="Insister sur les problèmes métier rencontrés avec l'approche précédente.")
    # Slide 4 - Pourquoi on change
    add_bullets_slide(prs, "Pourquoi on change (motifs métier)", [
        "Conversion et validation PDF/A intégrées → conformité d'archivage.",
        "Fiabilité : moins de documents rejetés lors d'export/archivage.",
        "Traçabilité : meilleurs rapports d'événements (warnings / erreurs).",
        "Moins de maintenance opérationnelle (SDK supporté).",
        "Support natif image→PDF/A pour normaliser les scans/photos."
    ], notes="Lier chaque motif à un bénéfice concret pour les produits/processus.")
    # Slide 5 - Fonctionnalités exposées
    add_bullets_slide(prs, "Fonctionnalités exposées (vue métier)", [
        "Conversion PDF → PDF/A (différents niveaux de conformité).",
        "Conversion Image → PDF/A (alternate text, langue).",
        "Fusion de documents en un seul PDF.",
        "Extraction de texte depuis un PDF.",
        "API orientée byte[] pour intégration aux services.",
        "Gestion automatique de fichiers temporaires (création & nettoyage)."
    ], notes="Donner exemples concrets : archivage d'un dossier, réception de scans.")
    # Slide 6 - Bénéfices métiers
    add_bullets_slide(prs, "Bénéfices métiers concrets", [
        "Documents archivés conformes → réduction du risque légal.",
        "Moins de retours/erreurs sur traitements batch.",
        "Standardisation du format de sortie (PDF/A).",
        "Meilleure visibilité en cas d'anomalie.",
        "Automatisation renforcée du traitement des pièces."
    ], notes="Aligner sur les KPIs métiers (p.ex. taux de rejet, temps de traitement).")
    # Slide 7 - Impacts (technique & infra)
    add_bullets_slide(prs, "Impacts à prévoir — technique & infra", [
        "Licence SDK à gérer et sécuriser (clé d'activation).",
        "Espace disque & permission d'écriture pour fichiers temporaires.",
        "Tests de montée en charge (I/O & mémoire).",
        "Dépendance binaire à déployer (.NET Framework 4.8).",
        "Monitoring des conversions (warnings / erreurs)."
    ], notes="Demander support DSI pour vault/secret manager et droits disque.")
    # Slide 8 - Impacts organisationnels & process
    add_bullets_slide(prs, "Impacts organisationnels & process", [
        "Procédure pour stocker/rotater la clé (ne pas laisser en clair).",
        "Clés distinctes pour dev/preprod/prod.",
        "QA métier : fournir jeux de documents pour validation.",
        "Support/exploitation : procédure d'escalade en cas d'échec."
    ], notes="Suggérer qui doit être responsable pour chaque point.")
    # Slide 9 - Risques et mitigations
    add_bullets_slide(prs, "Principaux risques (métier) & mitigations", [
        "Conversion non conforme pour certains fichiers → tests & validation PDF/A.",
        "Indisponibilité licence → clé de secours & contrat fournisseur.",
        "Problèmes de performance → tests de charge et montée en charge progressive.",
        "Fuite de clé → stocker dans vault et accès restreint."
    ], notes="Proposer un résultat souhaité après mitigation (ex. taux d'échec < X%).")
    # Slide 10 - Plan de migration
    add_bullets_slide(prs, "Plan de migration — étapes haut niveau", [
        "Préparation : inventaire usages & obtenir clés/licences (1-2 sem).",
        "Intégration & tests en dev (1-2 sem).",
        "Validation métier & QA (1-2 sem).",
        "Déploiement progressif (préprod → prod) et monitoring.",
        "Retrait d'iTextSharp après stabilisation."
    ], notes="Durée indicative totale : 2–6 semaines selon périmètre.")
    # Slide 11 - Critères d'acceptation
    add_bullets_slide(prs, "Critères d'acceptation (PO)", [
        "Documents convertis validés PDF/A selon niveau choisi.",
        "Fusion & extraction respectent besoins métiers.",
        "Pas de régression visible pour utilisateurs.",
        "Logs clairs et taux d'erreur acceptable.",
        "Clé licence en place et sécurisée."
    ], notes="Demander aux PO de définir seuils concrets (p.ex. latence max).")
    # Slide 12 - Décisions à prendre
    add_bullets_slide(prs, "Décisions à prendre aujourd'hui", [
        "Quel niveau PDF/A retenir pour l'archivage (ex. PDF/A-2b) ?",
        "PO fournissent jeux de documents représentatifs.",
        "Où stocker la clé (vault interne, secret manager) et qui la gère ?",
        "Planning souhaité pour le pilote et mise en production."
    ], notes="Fermer la réunion avec un ou deux décisions claires.")
    # Slide 13 - Démo / Scénarios recommandés
    add_bullets_slide(prs, "Démo — scénarios recommandés pour validation", [
        "Conversion d'un PDF non-conforme → vérifier PDF/A produit.",
        "Conversion d'un scan/image → vérifier PDF/A + texte alternatif.",
        "Fusion de 3 documents hétérogènes → vérifier pagination.",
        "Extraction de texte → vérifier qualité / encodage."
    ], notes="Proposer que les PO apportent 3-5 documents réels pour la session de validation.")
    # Slide 14 - Synthèse exécutive
    add_bullets_slide(prs, "Synthèse exécutive", [
        "Pourquoi : conformité d'archivage, fiabilité, réduction dette technique.",
        "Bénéfice : PDF/A natif, conversions images → archivables.",
        "Coûts/efforts : gestion licence, tests, adaptation infra.",
        "Proposition : lancer pilote rapide (dev → préprod) avec jeux fournis par PO."
    ], notes="Terminer sur appel à action clair.")
    # Slide 15 - Prochaines étapes
    add_bullets_slide(prs, "Prochaines étapes proposées", [
        "Validation du niveau PDF/A souhaité (PO).",
        "PO fournissent 3–5 documents types pour tests.",
        "Planifier le pilote (équipe technique) — ~2 sem.",
        "Réunion de revue des résultats & décision mise en prod."
    ], notes="Attribuer responsabilités et dates cibles.")
    # Slide 16 - Annexe : Checklist pour les PO
    add_bullets_slide(prs, "Annexe — Checklist rapide pour les PO", [
        "Choix du niveau PDF/A (ex. 2b).",
        "Fournir jeux de documents représentatifs (PDF, scans, gros fichiers).",
        "Définir seuils d'acceptation (latence, taux d'erreur).",
        "Confirmer responsable gestion de la clé licence."
    ], notes="Donner la checklist aux PO après la réunion.")
    # Save
    out_name = "Presentation_ApgisPdfTools_PO.pptx"
    prs.save(out_name)
    print(f"Présentation générée : {out_name}")

if __name__ == '__main__':
    main()
